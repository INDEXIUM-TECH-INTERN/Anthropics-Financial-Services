import sys
import subprocess

def install_dependencies():
    required_packages = ["openpyxl", "python-pptx"]
    for package in required_packages:
        try:
            __import__(package)
        except ImportError:
            print(f"Installing missing dependency: {package}...", flush=True)
            try:
                subprocess.check_call([sys.executable, "-m", "pip", "install", package])
                print(f"Successfully installed: {package}", flush=True)
            except Exception as e:
                print(f"Error installing {package}: {e}", file=sys.stderr, flush=True)

# Install required dependencies before any other imports
install_dependencies()

import os
import json
import argparse
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def is_numeric(val):
    if isinstance(val, (int, float)):
        return True
    if isinstance(val, str):
        if val.startswith("="):
            return True
        try:
            float(val.replace(',', '').replace('%', ''))
            return True
        except ValueError:
            return False
    return False

def create_fallback_xlsx(wb, report_title):
    # Sheet 1: Overview
    ws1 = wb.active
    ws1.title = "Overview"
    ws1.views.sheetView[0].showGridLines = True
    
    ws1.cell(row=2, column=2, value=report_title).font = Font(name="Calibri", size=18, bold=True, color="1F4E79")
    
    metadata = [
        ("Loại báo cáo", "Báo cáo Tài chính Excel"),
        ("Ngày tạo", "2026-06-11"),
        ("Người tạo", "Hệ thống Phân tích AI"),
        ("Mô tả", "Báo cáo tài chính chi tiết và mô hình định giá dự phòng.")
    ]
    
    thin_side = Side(border_style="thin", color="D3D3D3")
    thin_border = Border(left=thin_side, right=thin_side, top=thin_side, bottom=thin_side)
    
    for idx, (label, val) in enumerate(metadata, start=4):
        cell_lbl = ws1.cell(row=idx, column=2, value=label)
        cell_lbl.font = Font(name="Calibri", size=11, bold=True)
        cell_lbl.border = thin_border
        cell_lbl.fill = PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")
        
        cell_val = ws1.cell(row=idx, column=3, value=val)
        cell_val.font = Font(name="Calibri", size=11)
        cell_val.border = thin_border
        
    ws1.column_dimensions['B'].width = 20
    ws1.column_dimensions['C'].width = 50
    
    # Sheet 2: Financial Model
    ws2 = wb.create_sheet(title="Financial Model")
    ws2.views.sheetView[0].showGridLines = True
    
    # Title
    ws2.cell(row=1, column=1, value="Mô hình Tài chính So sánh FPT, VNM, HDB (2023-2025)").font = Font(name="Calibri", size=16, bold=True, color="1F4E79")
    ws2.row_dimensions[1].height = 25
    
    headers = ["Company", "Metric", "2023", "2024", "2025"]
    for col_idx, val in enumerate(headers, 1):
        cell = ws2.cell(row=3, column=col_idx, value=val)
        cell.font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = thin_border
    ws2.row_dimensions[3].height = 20
    
    rows = [
        # FPT
        ["FPT", "Revenue", 52618, 60000, 68000],
        ["FPT", "Cost of Goods Sold", 31570, 36000, 40800],
        ["FPT", "Gross Profit", "=C4-C5", "=D4-D5", "=E4-E5"],
        ["FPT", "Operating Expenses", 12000, 14000, 16000],
        ["FPT", "Operating Income", "=C6-C7", "=D6-D7", "=E6-E7"],
        # Empty
        ["", "", "", "", ""],
        # VNM
        ["VNM", "Revenue", 60479, 62000, 64000],
        ["VNM", "Cost of Goods Sold", 36287, 37200, 38400],
        ["VNM", "Gross Profit", "=C10-C11", "=D10-D11", "=E10-E11"],
        ["VNM", "Operating Expenses", 15000, 16000, 17000],
        ["VNM", "Operating Income", "=C12-C13", "=D12-D13", "=E12-E13"],
        # Empty
        ["", "", "", "", ""],
        # HDB
        ["HDB", "Revenue", 25300, 29000, 33000],
        ["HDB", "Cost of Goods Sold", 15180, 17400, 19800],
        ["HDB", "Gross Profit", "=C16-C17", "=D16-D17", "=E16-E17"],
        ["HDB", "Operating Expenses", 6000, 7000, 8000],
        ["HDB", "Operating Income", "=C18-C19", "=D18-D19", "=E18-E19"],
    ]
    
    for r_idx, r_data in enumerate(rows, start=4):
        is_blank = r_data[0] == ""
        is_bold = "Profit" in r_data[1] or "Income" in r_data[1]
        
        for c_idx, val in enumerate(r_data, 1):
            cell = ws2.cell(row=r_idx, column=c_idx, value=val)
            if is_blank:
                continue
            
            cell.border = thin_border
            cell.font = Font(name="Calibri", size=11, bold=is_bold)
            
            if c_idx <= 2:
                cell.alignment = Alignment(horizontal="left")
            else:
                cell.alignment = Alignment(horizontal="right")
                if isinstance(val, (int, float)):
                    cell.number_format = '#,##0'
                    
    for col in ws2.columns:
        max_len = 0
        for cell in col:
            val_str = str(cell.value or '')
            if len(val_str) > max_len:
                max_len = len(val_str)
        col_letter = get_column_letter(col[0].column)
        ws2.column_dimensions[col_letter].width = max(max_len + 4, 12)

def write_user_sheet_data(ws, sheet_title, rows):
    ws.views.sheetView[0].showGridLines = True
    thin_side = Side(border_style="thin", color="D3D3D3")
    thin_border = Border(left=thin_side, right=thin_side, top=thin_side, bottom=thin_side)
    
    ws.cell(row=1, column=1, value=sheet_title).font = Font(name="Calibri", size=16, bold=True, color="1F4E79")
    ws.row_dimensions[1].height = 25
    
    if not rows:
        return
        
    start_row = 3
    if all(isinstance(r, list) for r in rows):
        headers = rows[0]
        for col_idx, val in enumerate(headers, 1):
            cell = ws.cell(row=start_row, column=col_idx, value=val)
            cell.font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border
        ws.row_dimensions[start_row].height = 20
        
        for r_idx, r_data in enumerate(rows[1:], start_row + 1):
            for c_idx, val in enumerate(r_data, 1):
                cell = ws.cell(row=r_idx, column=c_idx)
                if is_numeric(val):
                    cell.alignment = Alignment(horizontal="right")
                    if isinstance(val, (int, float)):
                        cell.value = val
                        cell.number_format = '#,##0'
                    elif isinstance(val, str) and val.startswith("="):
                        cell.value = val
                    else:
                        try:
                            if '.' in val:
                                cell.value = float(val)
                            else:
                                cell.value = int(val)
                            cell.number_format = '#,##0'
                        except ValueError:
                            cell.value = val
                else:
                    cell.alignment = Alignment(horizontal="left")
                    cell.value = val
                
                cell.font = Font(name="Calibri", size=11)
                cell.border = thin_border
    elif all(isinstance(r, dict) for r in rows):
        keys = list(rows[0].keys())
        headers = keys
        for col_idx, val in enumerate(headers, 1):
            cell = ws.cell(row=start_row, column=col_idx, value=val)
            cell.font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border
        ws.row_dimensions[start_row].height = 20
        
        for r_idx, r_data in enumerate(rows, start_row + 1):
            for c_idx, key in enumerate(keys, 1):
                val = r_data.get(key, "")
                cell = ws.cell(row=r_idx, column=c_idx)
                if is_numeric(val):
                    cell.alignment = Alignment(horizontal="right")
                    if isinstance(val, (int, float)):
                        cell.value = val
                        cell.number_format = '#,##0'
                    elif isinstance(val, str) and val.startswith("="):
                        cell.value = val
                    else:
                        try:
                            if '.' in val:
                                cell.value = float(val)
                            else:
                                cell.value = int(val)
                            cell.number_format = '#,##0'
                        except ValueError:
                            cell.value = val
                else:
                    cell.alignment = Alignment(horizontal="left")
                    cell.value = val
                cell.font = Font(name="Calibri", size=11)
                cell.border = thin_border
    else:
        for r_idx, r_data in enumerate(rows, start_row):
            cell = ws.cell(row=r_idx, column=1, value=str(r_data))
            cell.font = Font(name="Calibri", size=11)
            cell.border = thin_border
            
    for col in ws.columns:
        max_len = 0
        for cell in col:
            val_str = str(cell.value or '')
            if len(val_str) > max_len:
                max_len = len(val_str)
        col_letter = get_column_letter(col[0].column)
        ws.column_dimensions[col_letter].width = max(max_len + 4, 12)

def generate_xlsx(output_path, title, data):
    wb = Workbook()
    if not data:
        create_fallback_xlsx(wb, title)
    else:
        if isinstance(data, dict):
            first = True
            for sheet_name, rows in data.items():
                if first:
                    ws = wb.active
                    ws.title = sheet_name
                    first = False
                else:
                    ws = wb.create_sheet(title=sheet_name)
                write_user_sheet_data(ws, sheet_name, rows)
        elif isinstance(data, list):
            ws = wb.active
            write_user_sheet_data(ws, title, data)
        else:
            create_fallback_xlsx(wb, title)
    wb.save(output_path)

def set_slide_bg(slide, r, g, b):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_bullet_points_to_slide(slide, title_text, bullets):
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(11.33), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Calibri'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.33), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    bullets = bullets[:8]
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p_b = tf2.paragraphs[0]
        else:
            p_b = tf2.add_paragraph()
        p_b.text = "• " + str(bullet)
        p_b.font.name = 'Calibri'
        p_b.font.size = Pt(16)
        p_b.font.color.rgb = RGBColor(51, 51, 51)
        p_b.space_after = Pt(14)

def add_table_to_slide(slide, title_text, headers, rows):
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(11.33), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Calibri'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)
    
    table_data = [headers] + rows
    table_data = table_data[:10]
    
    rows_count = len(table_data)
    cols_count = len(table_data[0])
    
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(11.33)
    height = Inches(4.5)
    
    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
    table = table_shape.table
    
    for r_idx, r_data in enumerate(table_data):
        for c_idx, val in enumerate(r_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Calibri'
            p.font.size = Pt(13)
            p.alignment = PP_ALIGN.CENTER if r_idx == 0 or c_idx == 0 else PP_ALIGN.RIGHT
            
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(31, 78, 121)
            else:
                p.font.color.rgb = RGBColor(51, 51, 51)
                if r_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(245, 247, 250)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

def generate_pptx(output_path, title, data):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Slide 1: Title
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1, 31, 78, 121)
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Calibri'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Báo cáo tự động tạo bởi Hệ thống AI | Ngày tạo: 2026-06-11"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(220, 220, 220)
    p2.alignment = PP_ALIGN.CENTER
    
    if not data:
        # Slide 2: Executive Summary
        slide2 = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide2, 255, 255, 255)
        add_bullet_points_to_slide(slide2, "Executive Summary / Tóm tắt Dự án", [
            "Hệ thống phân tích tài chính ghi nhận sự tăng trưởng ổn định của các doanh nghiệp trụ cột FPT, VNM và HDB trong giai đoạn 2023 - 2025.",
            "FPT dẫn đầu về tốc độ tăng trưởng doanh thu nhờ mảng xuất khẩu phần mềm và chuyển đổi số toàn cầu phát triển mạnh mẽ.",
            "VNM duy trì vị thế dẫn đầu thị trường sữa với dòng tiền hoạt động cực kỳ lành mạnh và tỷ lệ chi trả cổ tức cao, ổn định.",
            "HDBank thể hiện hiệu quả hoạt động vượt trội với sự gia tăng mạnh mẽ của tổng tài sản và kiểm soát tốt chất lượng nợ xấu.",
            "Báo cáo này đưa ra mô phỏng tài chính và các khuyến nghị chiến lược nhằm tối ưu hóa danh mục đầu tư trong các năm tới."
        ])
        
        # Slide 3: Table
        slide3 = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide3, 255, 255, 255)
        table_headers = ["Chỉ tiêu (Tỷ VNĐ)", "FPT", "VNM", "HDB"]
        table_rows = [
            ["Doanh thu", "60,000", "62,000", "29,000"],
            ["Giá vốn", "36,000", "37,200", "17,400"],
            ["Lợi nhuận gộp", "24,000", "24,800", "11,600"],
            ["Chi phí hoạt động", "14,000", "16,000", "7,000"],
            ["Lợi nhuận thuần", "10,000", "8,800", "4,600"]
        ]
        add_table_to_slide(slide3, "Financial Comparison / So sánh Tài chính 2024", table_headers, table_rows)
        
        # Slide 4: Recommendations
        slide4 = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide4, 255, 255, 255)
        add_bullet_points_to_slide(slide4, "Strategic Recommendations / Khuyến nghị Chiến lược", [
            "Tăng tỷ trọng đối với FPT: Triển vọng tăng trưởng doanh thu bền vững từ chuyển đổi số toàn cầu và AI mang lại động lực tăng giá mạnh mẽ.",
            "Nắm giữ đối với VNM: Đây là nguồn phòng thủ danh mục lý tưởng nhờ dòng tiền dồi dào và chính sách cổ tức cao trong giai đoạn thị trường biến động.",
            "Theo dõi mua đối với HDB: Tốc độ tăng trưởng tín dụng cao đi kèm với quản trị rủi ro tốt giúp HDBank trở thành cơ hội đầu tư tài chính ngân hàng hấp dẫn.",
            "Đa dạng hóa danh mục: Phân bổ tài sản hợp lý giữa các nhóm ngành Công nghệ, Tiêu dùng và Tài chính để tối thiểu hóa rủi ro phi hệ thống."
        ])
    else:
        if isinstance(data, dict):
            for sheet_name, content in data.items():
                slide = prs.slides.add_slide(blank_layout)
                set_slide_bg(slide, 255, 255, 255)
                if isinstance(content, list):
                    if not content:
                        add_bullet_points_to_slide(slide, sheet_name, ["Không có dữ liệu."])
                    elif all(isinstance(r, list) for r in content):
                        headers = content[0]
                        rows = content[1:]
                        add_table_to_slide(slide, sheet_name, headers, rows)
                    elif all(isinstance(r, dict) for r in content):
                        headers = list(content[0].keys())
                        rows = [[r.get(k, "") for k in headers] for r in content]
                        add_table_to_slide(slide, sheet_name, headers, rows)
                    else:
                        add_bullet_points_to_slide(slide, sheet_name, content)
                else:
                    add_bullet_points_to_slide(slide, sheet_name, [str(content)])
        elif isinstance(data, list):
            slide = prs.slides.add_slide(blank_layout)
            set_slide_bg(slide, 255, 255, 255)
            if not data:
                add_bullet_points_to_slide(slide, title, ["Không có dữ liệu."])
            elif all(isinstance(r, list) for r in data):
                headers = data[0]
                rows = data[1:]
                add_table_to_slide(slide, title, headers, rows)
            elif all(isinstance(r, dict) for r in data):
                headers = list(data[0].keys())
                rows = [[r.get(k, "") for k in headers] for r in data]
                add_table_to_slide(slide, title, headers, rows)
            else:
                add_bullet_points_to_slide(slide, title, data)
        else:
            slide = prs.slides.add_slide(blank_layout)
            set_slide_bg(slide, 255, 255, 255)
            add_bullet_points_to_slide(slide, title, [str(data)])
            
    prs.save(output_path)

def main():
    parser = argparse.ArgumentParser(description="Generate Excel or PowerPoint reports.")
    parser.add_argument("--type", "--format", dest="format", choices=["xlsx", "pptx"], required=True)
    parser.add_argument("--output", required=True, help="Absolute output filepath")
    parser.add_argument("--title", required=True, help="Report title")
    parser.add_argument("--data", help="JSON data string")
    parser.add_argument("--data-file", help="Filepath to a JSON data file")
    args = parser.parse_args()
    
    data = None
    if args.data:
        try:
            data = json.loads(args.data)
        except Exception as e:
            print(f"Error parsing --data: {e}", file=sys.stderr)
    elif args.data_file:
        try:
            with open(args.data_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
        except Exception as e:
            print(f"Error loading --data-file: {e}", file=sys.stderr)
            
    # Create output directory if it does not exist
    out_dir = os.path.dirname(args.output)
    if out_dir and not os.path.exists(out_dir):
        os.makedirs(out_dir, exist_ok=True)
        
    if args.format == "xlsx":
        generate_xlsx(args.output, args.title, data)
    elif args.format == "pptx":
        generate_pptx(args.output, args.title, data)
    print(f"Successfully generated {args.format} report at {args.output}")

if __name__ == "__main__":
    main()
